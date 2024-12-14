
import random
import openpyxl
from openpyxl.utils import get_column_letter
import datetime
import pptx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE  # Class in which the shape type is defined
from pptx.enum.shapes import MSO_SHAPE_TYPE
import numpy as np
from pptx.enum.text import PP_ALIGN
import easygui
from tkinter import filedialog
from tkinter import Text
from tkinter import ttk
from tkinter import *
from tkinter.filedialog import asksaveasfile
import tkinter as tk
from openpyxl.utils.exceptions import *
import easygui
import json
import os
import calendar
from enum import Enum


excelEventsDictionary = {}

def addToExcelEventDictionary(excelEvent, Dictionary):
    eventDayInMonth = excelEvent.dayInMonth
    if eventDayInMonth in Dictionary:
        Dictionary[eventDayInMonth].append(excelEvent)
        if len(Dictionary[eventDayInMonth]) > 2:
            raise Exception("בתאריך " + excelEvent.date + " קיימים מעל לשני אירועים") #TODO: check this case
    else:
        Dictionary[eventDayInMonth] = [excelEvent]


MAX_FIRST_MONTH_EVENT_COUNT = 31
MAX_SECOND_MONTH_EVENT_COUNT = 31
EXCEL_COULUMN_NUMBER = 9
FORBIDDEN_SENTENCE_HEBREW = "עלות פלוס: חינם | יאללה: חינם"
FORBIDDEN_SENTENCE_ARABIC = "السعر: شيكل"


class MetaData:
    def __init__(self, firstMonthInteger = None, secondMonthInteger = None, firstMonthName = None, secondMonthName = None, zone = None, year = None, contactName = None, contactPhone = None,
        language = None, increaseYear = False, hasCacheFile = False, imagesDirectory = None, allImages = None, usedImages = None):
        self.firstMonthInteger = firstMonthInteger
        self.secondMonthInteger = secondMonthInteger
        self.firstMonthName = firstMonthName
        self.secondMonthName = secondMonthName
        self.zone = zone
        self.year=year
        self.contactName = contactName
        self.contactPhone = contactPhone
        self.language = language
        self.increaseYear = increaseYear
        self.hasCacheFile = hasCacheFile
        self.imagesDirectory = imagesDirectory
        self.allImages = allImages
        self.usedImages = usedImages




class FirstSlideShapes:
    def __init__(self, monthsShape = None, yearShape = None, month1Shape = None, month2Shape = None, Singles1TitleShape = None, Singles1TextShape = None, YOLO1TitleShape = None, YOLO1TextShape = None,
        Kultura1TitleShape = None, Kultura1TextShape = None, Tiula1TitleShape = None, Tiula1TextShape = None, Golders1TitleShape = None, Golders1TextShape = None, Women1TitleShape = None,
        Women1TextShape = None, Singles2TitleShape = None, Singles2TextShape = None, YOLO2TitleShape = None, YOLO2TextShape = None, Kultura2TitleShape = None, Kultura2TextShape = None,
        Tiula2TitleShape = None, Tiula2TextShape = None, Golders2TitleShape = None, Golders2TextShape = None, Women2TitleShape = None, Women2TextShape = None, contactShape = None):
        self.monthsShape = monthsShape
        self.yearShape = yearShape
        self.month1Shape = month1Shape
        self.month2Shape = month2Shape
        self.Singles1TitleShape = Singles1TitleShape
        self.Singles1TextShape = Singles1TextShape
        self.YOLO1TitleShape = YOLO1TitleShape
        self.YOLO1TextShape = YOLO1TextShape
        self.Kultura1TitleShape = Kultura1TitleShape
        self.Kultura1TextShape = Kultura1TextShape
        self.Tiula1TitleShape = Tiula1TitleShape
        self.Tiula1TextShape = Tiula1TextShape
        self.Golders1TitleShape = Golders1TitleShape
        self.Golders1TextShape = Golders1TextShape
        self.Women1TitleShape = Women1TitleShape
        self.Women1TextShape = Women1TextShape
        self.Singles2TitleShape = Singles2TitleShape
        self.Singles2TextShape = Singles2TextShape
        self.YOLO2TitleShape = YOLO2TitleShape
        self.YOLO2TextShape = YOLO2TextShape
        self.Kultura2TitleShape = Kultura2TitleShape
        self.Kultura2TextShape = Kultura2TextShape
        self.Tiula2TitleShape = Tiula2TitleShape
        self.Tiula2TextShape = Tiula2TextShape
        self.Golders2TitleShape = Golders2TitleShape
        self.Golders2TextShape = Golders2TextShape
        self.Women2TitleShape = Women2TitleShape
        self.Women2TextShape = Women2TextShape
        self.contactShape = contactShape




# An object holding the fields of the Excel input event: date (datetime), day (string), hour (datetime.time), title (string), 
# text1 (string), location (string), price (string), isFree (bool), isZoom (bool) and month (int) of a planned event 
class ExcelEvent:
    def __init__(self, date, hour, title, location, price, community, link, eventType):
        self.date = date
        self.day = ""
        self.hour = hour
        self.title = title
        self.location = location
        self.price = price
        self.community = community
        self.link = link
        self.eventType = eventType
        self.month = getMonth(date)
        self.dayInMonth = getDay(date)


class Community(Enum):
    SINGLES = 1
    TIULA = 2
    YOLO = 3
    WOMEN = 4
    GOLDERS = 5
    KULTURA = 6
    YUMMIES = 7

class ExcelCommunitiesStrings:
    def __init__(self):
        self.Singles = "סינגלס"
        self.Tiula = "טיולא"
        self.Yolo = "YOLO"
        self.Women = "נשים"
        self.Golders = "גולדרס"
        self.Kultura = "קולטורה"
        self.Yummies = "יאמיס"
        self.excelCommunitiesStringsArray = [item[1] for item in vars(self).items() if not (item[0] == "excelCommunitiesStringsArray" or item == "communitiesStringsArray")]
        self.communitiesStringsArray = [item[0] for item in vars(self).items() if not (item[0] == "excelCommunitiesStringsArray" or item == "communitiesStringsArray")]

class EventType(Enum):
    REGULAR = 1
    WIDE = 2
    PEAK = 3

class EventTypes:
    def __init__(self):
        self.wideString = "אירוע מרחבי"
        self.wide = EventType.WIDE
        self.peakString = "אירוע שיא"
        self.peak = EventType.PEAK

class ShapeType(Enum):
    SINGLE = 1
    DOUBLE = 2
    PIC = 3
    OFF = 4


    

# An object holding the different shapes of the input powerpoint event shape: dateShape (shape), dayShape (shape), hourShape (shape), titleShape (shape), 
# text1Shape (shape), locationShape (shape) priceShape, freeShape (shape), zoomShape (shape)
# and shape (the root element holding them all together)
class SingleEventShape:
    def __init__(self, titleShape, locationShape, priceShape, tagSingles, tagTiula, tagYolo, tagWomen, tagGolders, tagKultura, tagYummies, countShape, countOffShape, dayShape, dayOffShape,
    spineBgShape, spineBgOffShape, bgHighlightShape, bgPicShape, bgOffShape, bgShape, shape):
        self.titleShape = titleShape
        self.locationShape = locationShape
        self.priceShape = priceShape
        self.tagSingles = tagSingles
        self.tagTiula = tagTiula
        self.tagYolo = tagYolo
        self.tagWomen = tagWomen
        self.tagGolders = tagGolders
        self.tagKultura = tagKultura
        self.tagYummies = tagYummies
        self.countShape = countShape
        self.countOffShape = countOffShape
        self.dayShape = dayShape
        self.dayOffShape = dayOffShape
        self.spineBgShape = spineBgShape
        self.spineBgOffShape = spineBgOffShape
        self.bgHighlightShape = bgHighlightShape
        self.bgPicShape = bgPicShape
        self.bgOffShape = bgOffShape
        self.bgShape = bgShape
        self.shape = shape

class DoubleEventShape:
    def __init__(self, titleShape1, locationShape1, priceShape1, tagSingles1, tagTiula1, tagYolo1, tagWomen1, tagGolders1, tagKultura1, tagYummies1,
     titleShape2, locationShape2, priceShape2, tagSingles2, tagTiula2, tagYolo2, tagWomen2, tagGolders2, tagKultura2, tagYummies2, countShape, countOffShape, dayShape, dayOffShape, 
     spineBgShape, spineBgOffShape, bgPicShape, bgOffShape, bgShape, shape):
        self.titleShape1 = titleShape1
        self.locationShape1 = locationShape1
        self.priceShape1 = priceShape1
        self.tagSingles1 = tagSingles1
        self.tagTiula1 = tagTiula1
        self.tagYolo1 = tagYolo1
        self.tagWomen1 = tagWomen1
        self.tagGolders1 = tagGolders1
        self.tagKultura1 = tagKultura1
        self.tagYummies1 = tagYummies1
        self.titleShape2 = titleShape2
        self.locationShape2 = locationShape2
        self.priceShape2 = priceShape2
        self.tagSingles2 = tagSingles2
        self.tagTiula2 = tagTiula2
        self.tagYolo2 = tagYolo2
        self.tagWomen2 = tagWomen2
        self.tagGolders2 = tagGolders2
        self.tagKultura2 = tagKultura2
        self.tagYummies2 = tagYummies2
        self.countShape = countShape
        self.countOffShape = countOffShape
        self.dayShape = dayShape
        self.dayOffShape = dayOffShape
        self.spineBgShape = spineBgShape
        self.spineBgOffShape = spineBgOffShape
        self.bgPicShape = bgPicShape
        self.bgOffShape = bgOffShape
        self.bgShape = bgShape
        self.shape = shape



# readExcel(excel_name)
#     Reading and analyzing the given excel file and returns a lists of events it contains
#     Arguments:
#         excel_name: a string path to the xlsx file holding
#             the data of all the planned events to be analyzed
#     Returns:
#         first_months_events: Event list containing all the events from excel that happen in the first month, 
#                              ignoring every event after MAX_FIRST_MONTH_EVENTS_COUNT
#         second_months_events: Event list containing all the events from excel that happen in the second month
#                              ignoring every event after MAX_SECOND_MONTH_EVENTs_COUNT
#         months: list of int containing the two months the excel deals with
#     Exceptions:
#         * CellCoordinatesException
#         * IllegalCharacterError
#         * InvalidFileException
#         * SheetTitleException
#         * MonthCountException: Less than two different months found in excel file
#         * MonthCountException: More than two different months found in excel file
#         * ValueError: time data does not match format

def readExcel(excel_name):

    # Define variable to load the dataframe
    try:
        ws = openpyxl.load_workbook(excel_name)

        # Define variable to read sheet
        sheet = ws.active
    
        getEventsFromExcel(sheet)
          
    except CellCoordinatesException:
        raise Exception("שגיאת המרה בין ערך נומרי ל-A1-style")
    except IllegalCharacterError:
        raise Exception("קובץ האקסל מכיל תוים לא חוקייים")
    except InvalidFileException:
        raise Exception("שיגאה בעת נסיון פתיחת קובץ שאינו קובץ אקסל")


#Split the Event list into first month and second month events
def splitExcelEventsByMonths(events):
    first_months_events = []
    second_months_events = []

    for event in events:
        if event.month == MetaData.firstMonthInteger:
            first_months_events.append(event)
        elif event.month == MetaData.secondMonthInteger:
            second_months_events.append(event)
        else:
            raise Exception("קובץ האקסל מכיל מעל לשני חודשיים")
    return first_months_events, second_months_events


def trimLeadingZero(inputString):
    if inputString.startswith('0'):
        return str(inputString[1:])
    else:
        return inputString

def checkDateStringValidity(inputString):
    if inputString != None:
        if len(inputString)==5:
            inputString = str(trimLeadingZero(inputString))
            if '.' in inputString:
                arrayString = inputString.split('.')
                if len(arrayString)==2:
                    if len(arrayString[1])==2:
                        return True
    return False

def getMonth(eventDate):
    eventDateArr = eventDate.split('.')
    if len(eventDateArr) < 2:
       eventDateMiddleSplit = eventDate.split(' ')
       eventDateArr = eventDateMiddleSplit[0].split('-')
    return int(eventDateArr[1])

def getDay(eventDate):
    eventDateArr = eventDate.split('.')
    if len(eventDateArr) < 2:
       eventDateMiddleSplit = eventDate.split(' ')
       eventDateArr = eventDateMiddleSplit[0].split('-')
    return int(eventDateArr[0])



# getEventsAndMonthsFromExcel(sheet)
#     Itereates through excel rows and cells to create Event and months lists
#     Arguments:
#         sheet: a Worksheet object to itereate through
#     Returns:
#         events: a sorted Event list containing all the events from the Worksheet
#         months: a sorted int list containing the months the excel Worksheet deals with
#     Exceptions:
#         ValueError

def getEventsFromExcel(sheet):
    excelMonths = []
    dictionary = firstMonthExcelEventsDictionary
    peakEventsDictionary = firstMonthPeakExcelEventsDictionary
    # Iterate the loop to read the cell values
    # Ignore first row with column titles
    for rowIndex in range(2, sheet.max_row+1):

        event_date = sheet.cell(row=rowIndex,column=1).value
        if not checkDateStringValidity(str(event_date)):
            if str(event_date)=="":
                if rowIsEmpty(sheet, rowIndex, EXCEL_COULUMN_NUMBER):
                    continue;
                else:
                    raise Exception("שורה \n"+str(rowIndex)+"\n עמודה: \n"+str(1)+"\nהערך ריק או לא תואם את התבנית:\n "+"DD.MM")

        day = sheet.cell(row=rowIndex,column=2).value
        if day is None: 
            raise Exception("שורה \n"+str(rowIndex)+"\n עמודה: \n"+str(2)+"\nהערך ריק ")

        hourValue = sheet.cell(row=rowIndex,column=3).value

        title = sheet.cell(row=rowIndex,column=4).value
        if title is None:
            title = ""

        location = sheet.cell(row=rowIndex,column=6).value
        if location is None:
            location = ""

        price = sheet.cell(row=rowIndex, column=7).value
        if price is None:
            price = ""

        communityString = sheet.cell(row = rowIndex, column = 10).value
        community = getCommunityFromString(communityString)

        link = sheet.cell(row = rowIndex, column = 11).value
        #TODO: check link validity

        eventTypeString = sheet.cell(row = rowIndex, column = 12).value
        eventType = getEventTypeFromString(eventTypeString)


        excelEvent = ExcelEvent(str(event_date), str(hourValue), title, location, price, community, link, eventType)

        if excelEvent.month not in excelMonths:
            excelMonths.append(excelEvent.month)
            if len(excelMonths) == 2:
                dictionary = secondMonthExcelEventsDictionary
                peakEventsDictionary = secondMonthPeakExcelEventsDictionary
            elif len(excelMonths) > 2:
                raise Exception("קובץ האקסל מכיל מעל לשני חודשים")

        if excelEvent.eventType is EventType.PEAK:
            peakEventsDictionary[excelEvent.community] = excelEvent
        else:
            addToExcelEventDictionary(excelEvent, dictionary)
             

    if len(excelMonths) < 2:
        raise Exception("קובץ האקסל מכיל פחות משני חודשים")

    if not (12 in excelMonths and 1 in excelMonths):
        excelMonths.sort()
        MetaData.increaseYear = False
    else:
        MetaData.increaseYear = True

    MetaData.firstMonthInteger = excelMonths[0]
    MetaData.secondMonthInteger = excelMonths[1]

    excelEventsDictionary [MetaData.firstMonthInteger] = firstMonthExcelEventsDictionary
    excelEventsDictionary[MetaData.secondMonthInteger] = secondMonthExcelEventsDictionary



def rowIsEmpty(sheet, rowIndex, maxColIndex):
    for i in range(1,maxColIndex+1):
        if sheet.cell(row=rowIndex,column=i).value !=None:
            return False
        else:
            if sheet.cell(row = rowIndex, column = 7) == None and sheet.cell(row=rowIndex, column =8) == False:
                return False
            return True


def getCommunityFromString(communityString):
    excelCommunitiesStrings = ExcelCommunitiesStrings()
    if excelCommunitiesStrings.Singles in communityString:
        return Community.SINGLES
    elif excelCommunitiesStrings.Tiula in communityString:
        return Community.TIULA
    elif excelCommunitiesStrings.Yolo in communityString:
        return Community.YOLO
    elif excelCommunitiesStrings.Women in communityString:
        return Community.WOMEN
    elif excelCommunitiesStrings.Golders in communityString:
        return Community.GOLDERS
    elif excelCommunitiesStrings.Kultura in communityString:
        return Community.KULTURA
    elif excelCommunitiesStrings.Yummies in communityString:
        return Community.YUMMIES
    else:
        return None


def getEventTypeFromString(eventTypeString):
    eventTypes = EventTypes()
    if eventTypes.wideString in eventTypeString:
        return EventType.WIDE
    elif eventTypes.peakString in eventTypeString:
        return EventType.PEAK
    else:
        return EventType.REGULAR

 
# get_text_boxes(slide, months)
#     Itereates through all powerpoint shapes (using iter_textframed_shapes method) that contain textboxes, sort them,
#     create month textboxes and returns an event_shapes sorted list of all textable leaf (not group) EventShape objects
#     Arguments:
#         slide: a Presentation Slide object to analize for shapes
#         months: array of ints, representing the numeric (int) representation of months presenting in the excel events sheet
#     Returns:
#         event_shapes: a sorted list of all textable leaf (not group) EventShape objects in the given slide
#     Exceptions:
#         * Not all cells of Excel are filled

def get_slide_shapes(slide, isFirstSlide = False):
    singles, doubles, header_shape = find_groups(slide.shapes, isFirstSlide)[:3]
    first_slide_groups, first_slide_month1_shape, first_slide_month2_shape, first_slide_phone_shape = find_groups(slide.shapes, isFirstSlide)[3:]
    header_shape_dict = {}
    first_slide_shapes_dict = {}

    if not isFirstSlide:

        single_titles = []
        single_locations = []
        single_prices = [] 
        single_Singles_tags = []
        single_Tiula_tags = []
        single_Yolo_tags = []
        single_Women_tags = []
        single_Golders_tags = []
        single_Kultura_tags = []
        single_Yummies_tags = []
        single_counts = []
        single_count_offs = []
        single_days = []
        single_day_offs = []
        single_spine_bgs = []
        single_spine_bg_offs = []
        single_bg_highlights = []
        single_bg_pics = []
        single_bg_offs = []
        single_bgs = []
        single_shapes = []

        double_titles1 = []
        double_locations1 = []
        double_prices1 = [] 
        double_Singles_tags1 = []
        double_Tiula_tags1 = []
        double_Yolo_tags1 = []
        double_Women_tags1 = []
        double_Golders_tags1 = []
        double_Kultura_tags1 = []
        double_Yummies_tags1 = []
        double_titles2 = []
        double_locations2 = []
        double_prices2 = [] 
        double_Singles_tags2 = []
        double_Tiula_tags2 = []
        double_Yolo_tags2 = []
        double_Women_tags2 = []
        double_Golders_tags2 = []
        double_Kultura_tags2 = []
        double_Yummies_tags2 = []
        double_counts = []
        double_count_offs = []
        double_days = []
        double_day_offs = []
        double_spine_bgs = []
        double_spine_bg_offs = []
        double_bg_pics = []
        double_bg_offs = []
        double_bgs = []
        double_shapes = []

        for shape in header_shape.shapes:
            if shape.name == 'MONTH':
                header_shape_dict["month"] = shape
            elif shape.name == 'YEAR':
                header_shape_dict["year"] = shape
            elif shape.name == 'ZONE':
                header_shape_dict["zone"] = shape

            
        for g in singles:
            single_shapes.append(g)
            for shape in g.shapes:
                if shape.name == 'TITLE':
                    single_titles.append(shape)
                elif shape.name == 'LOCATION':
                    single_locations.append(shape)
                elif shape.name == 'PRICE':
                    single_prices.append(shape)
                elif shape.name == 'TAG SINGLES 1':
                    single_Singles_tags.append(shape)
                elif shape.name == 'TAG TIULA 1':
                    single_Tiula_tags.append(shape)
                elif shape.name == 'TAG YOLO 1':
                    single_Yolo_tags.append(shape)
                elif shape.name == 'TAG WOMEN 1':
                    single_Women_tags.append(shape)
                elif shape.name == 'TAG GOLDERS 1':
                    single_Golders_tags.append(shape)
                elif shape.name == 'TAG KULTURA 1':
                    single_Kultura_tags.append(shape)
                elif shape.name == 'TAG YUMMIES 1':
                    single_Yummies_tags.append(shape)
                elif shape.name == 'COUNT':
                    single_counts.append(shape)
                elif shape.name == 'COUNT OFF':
                    single_count_offs.append(shape)
                elif shape.name == 'DAY':
                    single_days.append(shape)
                elif shape.name == 'DAY OFF':
                    single_day_offs.append(shape)
                elif shape.name == 'SPINE BG':
                    single_spine_bgs.append(shape)
                elif shape.name == 'SPINE BG OFF':
                    single_spine_bg_offs.append(shape)
                elif shape.name == 'BG HIGHLIGHT':
                    single_bg_highlights.append(shape)
                elif shape.name == 'BG PIC':
                    single_bg_pics.append(shape)
                elif shape.name == 'BG OFF':
                    single_bg_offs.append(shape)
                elif shape.name == 'BG':
                    single_bgs.append(shape)


        for g in doubles:
            double_shapes.append(g)
            for shape in g.shapes:
                if shape.name == 'TITLE 1':
                    double_titles1.append(shape)
                elif shape.name == 'LOCATION 1':
                    double_locations1.append(shape)
                elif shape.name == 'PRICE 1':
                    double_prices1.append(shape)
                elif shape.name == 'TAG SINGLES 1':
                    double_Singles_tags1.append(shape)
                elif shape.name == 'TAG TIULA 1':
                    double_Tiula_tags1.append(shape)
                elif shape.name == 'TAG YOLO 1':
                    double_Yolo_tags1.append(shape)
                elif shape.name == 'TAG WOMEN 1':
                    double_Women_tags1.append(shape)
                elif shape.name == 'TAG GOLDERS 1':
                    double_Golders_tags1.append(shape)
                elif shape.name == 'TAG KULTURA 1':
                    double_Kultura_tags1.append(shape)
                elif shape.name == 'TAG YUMMIES 1':
                    double_Yummies_tags1.append(shape)
                elif shape.name == 'TITLE 2':
                    double_titles2.append(shape)
                elif shape.name == 'LOCATION 2':
                    double_locations2.append(shape)
                elif shape.name == 'PRICE 2':
                    double_prices2.append(shape)
                elif shape.name == 'TAG SINGLES 2':
                    double_Singles_tags2.append(shape)
                elif shape.name == 'TAG TIULA 2':
                    double_Tiula_tags2.append(shape)
                elif shape.name == 'TAG YOLO 2':
                    double_Yolo_tags2.append(shape)
                elif shape.name == 'TAG WOMEN 2':
                    double_Women_tags2.append(shape)
                elif shape.name == 'TAG GOLDERS 2':
                    double_Golders_tags2.append(shape)
                elif shape.name == 'TAG KULTURA 2':
                    double_Kultura_tags2.append(shape)
                elif shape.name == 'YUMMIES 2':
                    double_Yummies_tags2.append(shape)
                elif shape.name == 'COUNT':
                    double_counts.append(shape)
                elif shape.name == 'COUNT OFF':
                    double_count_offs.append(shape)
                elif shape.name == 'DAY':
                    double_days.append(shape)
                elif shape.name == 'DAY OFF':
                    double_day_offs.append(shape)
                elif shape.name == 'SPINE BG':
                    double_spine_bgs.append(shape)
                elif shape.name == 'SPINE BG OFF':
                    double_spine_bg_offs.append(shape)
                elif shape.name == 'BG PIC':
                    double_bg_pics.append(shape)
                elif shape.name == 'BG OFF':
                    double_bg_offs.append(shape)
                elif shape.name == 'BG':
                    double_bgs.append(shape)

        single_event_shapes = []
        double_event_shapes = []

        for i in range(36):
            single_event_shape = SingleEventShape(single_titles[i], single_locations[i], single_prices[i], single_Singles_tags[i], single_Tiula_tags[i], single_Yolo_tags[i], single_Women_tags[i],
                single_Golders_tags[i], single_Kultura_tags[i], single_Yummies_tags[i], single_counts[i], single_count_offs[i], single_days[i], single_day_offs[i], single_spine_bgs[i], single_spine_bg_offs[i],
                single_bg_highlights[i], single_bg_pics[i], single_bg_offs[i], single_bgs[i], single_shapes[i])
            single_event_shapes.append(single_event_shape)

            double_event_shape = DoubleEventShape(double_titles1[i], double_locations1[i], double_prices1[i], double_Singles_tags1[i], double_Tiula_tags1[i], double_Yolo_tags1[i], double_Women_tags1[i],
                double_Golders_tags1[i], double_Kultura_tags1[i], double_Yummies_tags1[i], double_titles2[i], double_locations2[i], double_prices2[i], double_Singles_tags2[i], double_Tiula_tags2[i], double_Yolo_tags2[i],
                double_Women_tags2[i], double_Golders_tags2[i], double_Kultura_tags2[i],double_Yummies_tags2[i], double_counts[i], double_count_offs[i], double_days[i], double_day_offs[i], double_spine_bgs[i],
                double_spine_bg_offs[i], double_bg_pics[i], double_bg_offs[i], double_bgs[i], double_shapes[i])
            double_event_shapes.append(double_event_shape)

        return single_event_shapes, double_event_shapes, header_shape_dict
    
    else:

        for g in first_slide_groups:
            if g.name == "SINGLES 1":
                for shape in g.shapes:
                    if shape.name == "title":
                        first_slide_shapes_dict["Singles1TitleShape"] = shape
                    elif shape.name == "info":
                        first_slide_shapes_dict["Singles1TextShape"] = shape
            elif g.name == "SINGLES 2":
                for shape in g.shapes:
                    if shape.name == "title":
                        first_slide_shapes_dict["Singles2TitleShape"] = shape
                    elif shape.name == "info":
                        first_slide_shapes_dict["Singles2TextShape"] = shape
            elif g.name == "YOLO 1":
                for shape in g.shapes:
                    if shape.name == "title":
                        first_slide_shapes_dict["Yolo1TitleShape"] = shape
                    elif shape.name == "info":
                        first_slide_shapes_dict["Yolo1TextShape"] = shape
            elif g.name == "YOLO 2":
                for shape in g.shapes:
                    if shape.name == "title":
                        first_slide_shapes_dict["Yolo2TitleShape"] = shape
                    elif shape.name == "info":
                        first_slide_shapes_dict["Yolo2TextShape"] = shape
            elif g.name == "KULTURA 1":
                for shape in g.shapes:
                    if shape.name == "title":
                        first_slide_shapes_dict["Kultura1TitleShape"] = shape
                    elif shape.name == "info":
                        first_slide_shapes_dict["Kultura1TextShape"] = shape
            elif g.name == "KULTURA 2":
                for shape in g.shapes:
                    if shape.name == "title":
                        first_slide_shapes_dict["Kultura2TitleShape"] = shape
                    elif shape.name == "info":
                        first_slide_shapes_dict["Kultura2TextShape"] = shape
            elif g.name == "TIULA 1":
                for shape in g.shapes:
                    if shape.name == "title":
                        first_slide_shapes_dict["Tiula1TitleShape"] = shape
                    elif shape.name == "info":
                        first_slide_shapes_dict["Tiula1TextShape"] = shape
            elif g.name == "TIULA 2":
                for shape in g.shapes:
                    if shape.name == "title":
                        first_slide_shapes_dict["Tiula2TitleShape"] = shape
                    elif shape.name == "info":
                        first_slide_shapes_dict["Tiula2TextShape"] = shape
            elif g.name == "GOLDERS 1":
                for shape in g.shapes:
                    if shape.name == "title":
                        first_slide_shapes_dict["Golders1TitleShape"] = shape
                    elif shape.name == "info":
                        first_slide_shapes_dict["Golders1TextShape"] = shape
            elif g.name == "GOLDERS 2":
                for shape in g.shapes:
                    if shape.name == "title":
                        first_slide_shapes_dict["Golders2TitleShape"] = shape
                    elif shape.name == "info":
                        first_slide_shapes_dict["Golders2TextShape"] = shape
            elif g.name == "WOMEN 1":
                for shape in g.shapes:
                    if shape.name == "title":
                        first_slide_shapes_dict["Women1TitleShape"] = shape
                    elif shape.name == "info":
                        first_slide_shapes_dict["Women1TextShape"] = shape
            elif g.name == "WOMEN 2":
                for shape in g.shapes:
                    if shape.name == "title":
                        first_slide_shapes_dict["Women2TitleShape"] = shape
                    elif shape.name == "info":
                        first_slide_shapes_dict["Women2TextShape"] = shape
            elif g.name == "YUMMIES 1":
                for shape in g.shapes:
                    if shape.name == "title":
                        first_slide_shapes_dict["Yummies1TitleShape"] = shape
                    elif shape.name == "info":
                        first_slide_shapes_dict["Yummies1TextShape"] = shape
            elif g.name == "YUMMIES 2":
                for shape in g.shapes:
                    if shape.name == "title":
                        first_slide_shapes_dict["Yummies2TitleShape"] = shape
                    elif shape.name == "info":
                        first_slide_shapes_dict["Yummies2TextShape"] = shape
            elif g.name == "FOOTER":
                for shape in g.shapes:
                    if shape.name == "NAME AND TEL":
                        first_slide_shapes_dict["contact"] = shape
            elif g.name == "HEADER":
                for shape in g.shapes:
                    if shape.name == "MONTHES":
                        first_slide_shapes_dict["monthes"] = shape
                    elif shape.name == "YEAR":
                        first_slide_shapes_dict["year"] = shape

        first_slide_shapes_dict["month1"] = first_slide_month1_shape
        first_slide_shapes_dict["month2"] = first_slide_month2_shape
        first_slide_shapes_dict["phone"] = first_slide_phone_shape
        print(first_slide_shapes_dict["phone"])

        return first_slide_shapes_dict
        

def find_groups(shapes, isFirstSlide = True):

    doubles = []
    singles = []
    header = None
    first_slide_groups = []
    first_slide_month1_shape = None 
    first_slide_month2_shape = None
    first_slide_phone_shape = None

    if not isFirstSlide:

        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                if shape.name.startswith("DOUBLE"):
                    doubles.append(shape)
                elif shape.name.startswith("SINGLE"):
                    singles.append(shape)
                elif shape.name.startswith("HEADER"):
                    header = shape

    else:

        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                if shape.name.startswith("HEADER"):
                    first_slide_groups.append(shape)
                elif shape.name.startswith("SINGLES"):
                    first_slide_groups.append(shape)
                elif shape.name.startswith("YOLO"):
                    first_slide_groups.append(shape)
                elif shape.name.startswith("KULTURA"):
                    first_slide_groups.append(shape)  
                elif shape.name.startswith("TIULA"):
                    first_slide_groups.append(shape)
                elif shape.name.startswith("GOLDERS"):
                    first_slide_groups.append(shape)
                elif shape.name.startswith("WOMEN"):
                    first_slide_groups.append(shape)
                elif shape.name.startswith("YUMMIES"):
                    first_slide_groups.append(shape)
                elif shape.name.startswith("FOOTER"):
                    first_slide_groups.append(shape)
            else:
                if shape.name == "TITLE MONTH 1":
                    first_slide_month1_shape = shape
                elif shape.name == "TITLE MONTH 2":
                    first_slide_month2_shape = shape
                elif shape.name == "PHONE CONTACT":
                    first_slide_phone_shape = shape
                    print(first_slide_phone_shape)

    return singles, doubles, header, first_slide_groups, first_slide_month1_shape, first_slide_month2_shape, first_slide_phone_shape


# iter_textframed_shapes(shapes):
#     Itereates through all powerpoint shapes (using iter_textable_shapes method) that contain textboxes and yields
#     the leaf shapes with text boxes and their parent group shapes
#     Arguments:
#         shapes: a list of shapes (Shape objects) to iterate throuth
#     Returns:
#         a leaf shape with textbox or a group shape hilding shapes with textboxes

def iter_textframed_shapes(shapes):
    """Generate shape objects in shapes that can contain text.

    Shape objects are generated in document order (z-order), bottom to top.
    """
    for shape in shapes:
        # ---recurse on group shapes---
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            group_shape = shape
            yield shape
            iter_textframed_shapes(group_shape.shapes)
            for shape in group_shape.shapes:
                yield shape
            continue

        # ---otherwise, treat shape as a "leaf" shape---
        if shape.has_text_frame:
            yield shape


# iter_textable_shapes(shapes):
#     Itereates through shapes in group shape to find and yield the leaf shape with textboxes
#     Arguments:
#         shapes: a list of shapes (Shape objects) to iterate throuth
#     Returns:
#         a leaf shape with textbox

# def iter_textable_shapes(shapes):
#     for shape in shapes:
#             yield shape


# create_program_GUI():
#     creates the program GUI window

def create_program_GUI():

    useUserSavedInput = False
    userSavedInput = None

    try:
        cache_file = open("cache.txt", "r")
        content = cache_file.read()
        if os.stat("cache.txt").st_size != 0:
            userSavedInput = json.loads(content)
            useUserSavedInput = True
        cache_file.close()
    except FileNotFoundError:
        pass

    def assignExcelPath(value):
        global excelFilePath
        validity, extension = checkValidity(value, "xlsx")
        excelFilePath = value
        textbox1.configure(state="normal")
        if len(textbox1.get("1.0", tk.END))>1:
            textbox1.delete(1.0,tk.END)
        textbox1.insert(tk.END, excelFilePath)
        textbox1.configure(state="disabled")

    def assignPowerpointPath(value):
        global pptxFilePath
        pptxFilePath = value
        validity, extension = checkValidity(value, "pptx")
        textbox2.configure(state="normal")
        if len(textbox2.get("1.0", tk.END))>1:
            textbox2.delete(1.0,tk.END)
        textbox2.insert(tk.END, pptxFilePath)
        textbox2.configure(state="disabled")

    try:

        app = tk.Tk()
        app.wm_title("2MonthPlanCreator")
        app.geometry("650x250")
        app.padx = (100,0)
        app['background']='#015293'

        headLabel = tk.Label(text='פרטים לתוכנית ארועים',  font='Arial 14 bold')
        headLabel.grid(sticky="e", row = 0, column = 5, columnspan = 2)
        headLabel['background']='#015293'
        headLabel.config(fg= "#fdc53e")

        assignFirstMonthLabel = tk.Label(text="חודש ראשון")
        assignFirstMonthLabel.grid(sticky="e", row = 1, column =6)
        assignFirstMonthLabel['background']='#015293'
        assignFirstMonthLabel.config(fg="white")

        assignFirstMonthTextbox = ttk.Entry(justify=tk.RIGHT)
        assignFirstMonthTextbox.grid(stick="e", row=1, column=5)
        if useUserSavedInput == True and userSavedInput != None:
            assignFirstMonthTextbox.insert(0, userSavedInput.get("month1"))


        assignSecondMonthLabel = tk.Label(text="חודש שני")
        assignSecondMonthLabel.grid(sticky="e", row=1, column = 4)
        assignSecondMonthLabel['background']='#015293'
        assignSecondMonthLabel.config(fg="white")

        assignSecondMonthTextbox = ttk.Entry(justify=tk.RIGHT)
        assignSecondMonthTextbox.grid(sticky="e", row=1, column = 3, padx=(10,10))
        if useUserSavedInput == True and userSavedInput != None:
            assignSecondMonthTextbox.insert(0, userSavedInput.get("month2"))

        assignYearLabel = tk.Label(text = "שנה")
        assignYearLabel.grid(sticky="e", row = 1, column = 2)
        assignYearLabel['background']='#015293'
        assignYearLabel.config(fg="white")

        assignYearTextbox = ttk.Entry(justify=tk.RIGHT, width="10")
        assignYearTextbox.grid(sticky="e", row = 1, column = 1, padx=(10,10))
        if useUserSavedInput == True and userSavedInput != None:
            assignYearTextbox.insert(0, userSavedInput.get("year"))

        language = tk.StringVar()

        hebrew = Radiobutton(app, text="עברית", variable = language, value = 1)
        hebrew['background']='#015293'
        hebrew.config(fg="white", selectcolor="#4cb263")
        hebrew.grid(row = 2, column = 6, sticky="e")
        arabic = Radiobutton(app, text="ערבית", variable = language, value = 2)
        arabic['background']='#015293'
        arabic.config(fg="white", selectcolor="#4cb263")
        arabic.grid(row = 2, column = 5, sticky="e")

        if useUserSavedInput == True and userSavedInput != None:
            language.set(userSavedInput.get("language"))
        else:
            language.set(1)

        assignLocationLabel = tk.Label(text="מיקום")
        assignLocationLabel.grid(sticky="e", row = 3, column = 6)
        assignLocationLabel['background']='#015293'
        assignLocationLabel.config(fg="white")

        assignLocationTextbox = ttk.Entry(justify=tk.RIGHT)
        assignLocationTextbox.grid(sticky="e", row = 3, column = 5)
        if useUserSavedInput == True and userSavedInput != None:
            assignLocationTextbox.insert(0, userSavedInput.get("area"))

        assignContactLabel = tk.Label(text = "לקבלת פרטים לפנות אל")
        assignContactLabel.grid(sticky="e", row = 4, column = 6)
        assignContactLabel['background']='#015293'
        assignContactLabel.config(fg="white")

        assignContactTextbox = ttk.Entry(justify=tk.RIGHT)
        assignContactTextbox.grid(sticky="e", row = 4, column = 5)
        if useUserSavedInput == True and userSavedInput != None:
            assignContactTextbox.insert(tk.END, userSavedInput.get("contact"))

        directions1 = tk.Label(text = "אנא בחרו קובץ אקסל עם המידע הרלוונטי וקובץ פאוורפוינט לדוגמא", justify="right")
        directions1['background']='#015293'
        directions1.config(fg="white")
        directions1.grid(sticky="e", row=5, column=4, columnspan=4)

        btn1 = tk.Button(text="בחרו קובץ אקסל", command=lambda: assignExcelPath(filedialog.askopenfilename(filetypes=[("קובץ אקסל" , ".xlsx")])))
        btn1.grid(sticky="e", row=6, column=6)
        btn1['background']='#008fd1'
        btn1.config(fg="white")

        # Create text widget and specify size.
        textbox1 = Text(app, height = 1, width = 55)
        textbox1.grid(sticky="e", row=6, column=1, columnspan=5)
        textbox1.configure(state="disabled")

        btn2 = tk.Button(text="בחרו קובץ פאוורפוינט", command=lambda: assignPowerpointPath(filedialog.askopenfilename(filetypes=[("קובץ פאוורפוינט" , ".pptx")])))
        btn2.grid(sticky="e", row=7, column = 6)
        btn2['background']='#008fd1'
        btn2.config(fg="white")

        # Create text widget and specify size.
        textbox2 = Text(app, height = 1, width = 55)
        textbox2.grid(sticky="e", row=7, column=1, columnspan=5)
        textbox2.configure(state="disabled")

        btn4 = tk.Button(text="צרו תוכנית דו-חודשית!", command= lambda: createPptxPlans(assignFirstMonthTextbox.get(), assignSecondMonthTextbox.get(), assignLocationTextbox.get(), assignContactTextbox.get(),assignYearTextbox.get(), language.get()))
        btn4.grid(sticky="e", row=9, column = 5)
        btn4['background']='#4cb263'
        btn4.config(fg="white")


        btn5 = tk.Button(text="סגירה", command=app.destroy)
        btn5.grid(sticky="e", row=9, column = 4)
        btn5['background']='#008fd1'
        btn5.config(fg="white")

        app.mainloop()  

    except Exception as e:
        easygui.msgbox(e)
        raise Exception("קובץ ה-cache אינו תקין")




# checkValidity(value, extension):
#     Checks if the given filepath's extension matches the desired extesion and returns 
#     Arguments:
#         value: the given filepath to be checked
#         extension: the desired extension
#     Returns:
#         a tuple of: 
#             True/False depending on whether the given value file extension matches the desired extension
#             a string containing the real extension of the file given 
#             if the value does not contain an extension ('.xxx') returns (False, "")

def checkValidity(value, extension):
    if extension == "":
    #check for filename validity
        try:
            nameExtension = value.split('.')
            if len(nameExtension)>2:
                return False, ""
            elif len(nameExtension) == 2 and nameExtension[1]!="pptx":
                return True, ""
            elif len(nameExtension) ==2 and nameExtension[1] == "[pptx]":
                return True, ""
            else:
                return True, "pptx"
        except:
            return False, ""
    else:
        try:
            valueExtension = value.split('.')[1]
            valueExtension = ''.join(c for c in valueExtension if c.isprintable())
            return valueExtension == extension, valueExtension
        except:
            return False, ""


# createPptxPlan(month1, month2, area, contact, year, language):
#     Creates and saves the two-month plan pptx file 
#     Arguments:
#         month1:
#         month2:
#         area:
#         contact:
#         year:
#         language: 

def createPptxPlans(month1, month2, area, contact, year, language):
    resetGlobalVariables()
    if len(year)!=4:
            raise Exception("שדה השנה צריך להכיל 4 תוים בדיוק. למשל: 2023")
    try: 
        cache_file = open("cache.txt", "w")
        userInput = {
            "month1": month1,
            "month2": month2,
            "area": area,
            "contact": contact,
            "year": year,
            "language": language
        }
        cache_file.write(json.dumps(userInput))
        cache_file.close

        resolveMetaData(month1, month2, area, contact, year, language)

        readExcel(excelFilePath)

        try:

            presentation = Presentation(pptxFilePath)

            first_slide = presentation.slides[0]
            second_slide = presentation.slides[1]

            processFirstSlide(first_slide)
            processSecondSlide(second_slide, True)

            defaultFileName = "תוכניה " + month1 + "-" + month2

            files = [('Powerpoint files', '*.pptx')]
            saved_file = asksaveasfile(filetypes = files, defaultextension = files, initialfile  = defaultFileName + " - 1") #TODO: add default file name
            saved_file2 = asksaveasfile(filetypes = files, defaultextension = files, initialfile  = defaultFileName + " - 2")
            try:
                presentation.save(saved_file.name)
            except PermissionError: 
                easygui.msgbox("הקובץ \n"+ saved_file.name+"\n פתוח. יש לסגור אותו בעת הרצת התוכנית")

            presentation = Presentation(pptxFilePath)

            first_slide = presentation.slides[0]
            second_slide = presentation.slides[1]

            processFirstSlide(first_slide)
            processSecondSlide(second_slide, False)

            try:
                presentation.save(saved_file2.name)
                easygui.msgbox("הקבצים נוצרו בהצלחה!")
            except PermissionError: 
                easygui.msgbox("הקובץ \n"+ saved_file2.name+"\n פתוח. יש לסגור אותו בעת הרצת התוכנית")

        except IndexError as e:
            if len(presentation.slides) > 2:
                raise Exception("תבנית ה-powerpoint מכילה פחות מ-2 שקפים")
            else:
                raise Exception("שגיאה: Index Out of Range")

    except Exception as e:
        easygui.msgbox("שגיאה :"+ str(e))


def resolveMetaData(month1, month2, area, contact, year, language):
    MetaData.firstMonthName = month1
    MetaData.secondMonthName = month2
    MetaData.zone = area
    MetaData.contact = contact
    MetaData.year = year
    MetaData.language = language
    dir = MetaData.imagesDirectory = 'AppData/images'

    if not os.path.isdir(dir):
        raise ValueError(f"Directory {dir} does not exist")
    
    MetaData.allImages = [
        f for f in os.listdir(dir) 
        if os.path.isfile(os.path.join(dir, f)) 
        and os.path.splitext(f)[1].lower() == '.png'
    ]

    if not MetaData.allImages:
        raise ValueError(f"No image files found in {dir}")

    MetaData.usedImages = set()


def resetGlobalVariables():
    global firstMonthExcelEventsDictionary
    global secondMonthExcelEventsDictionary
    global firstMonthPeakExcelEventsDictionary
    global secondMonthPeakExcelEventsDictionary
    firstMonthExcelEventsDictionary = {}
    secondMonthExcelEventsDictionary = {}
    firstMonthPeakExcelEventsDictionary = {}
    secondMonthPeakExcelEventsDictionary = {}

    


def processFirstSlide(slide):
    first_slide_shapes_dict = get_slide_shapes(slide, isFirstSlide = True)
    month1 = MetaData.firstMonthName
    month2 = MetaData.secondMonthName
    year = MetaData.year
    year_string = year[2:]
    increaseYear = MetaData.increaseYear
    increased_year_string = str(int(year_string) + 1) if increaseYear else year_string
    contact = MetaData.contact

    excelCommunitiesStrings = ExcelCommunitiesStrings()
    dictionaries = [firstMonthPeakExcelEventsDictionary, secondMonthPeakExcelEventsDictionary]

    for index, commString in enumerate(excelCommunitiesStrings.communitiesStringsArray):
        community = Community[commString.upper()]
        for i in range(1,3):
            dictionary = dictionaries[i-1]
            titleKey = commString + str(i) + "TitleShape"
            textKey = commString + str(i) + "TextShape"
            if community in dictionary:
                full_date = create_date_string(dictionary[community].date, increased_year_string)
                title = excelCommunitiesStrings.excelCommunitiesStringsArray[index] + " - " + dictionary[community].title
                text = dictionary[community].location + " | " + full_date + " | " + dictionary[community].hour
                writeTextToTextbox(first_slide_shapes_dict[titleKey].text_frame, title, link = dictionary[community].link)
                writeTextToTextbox(first_slide_shapes_dict[textKey].text_frame, text)
            else:
                title = excelCommunitiesStrings.excelCommunitiesStringsArray[index]
                text = "ניפגש בחודשים הבאים!"
                writeTextToTextbox(first_slide_shapes_dict[titleKey].text_frame, title)
                writeTextToTextbox(first_slide_shapes_dict[textKey].text_frame, text)

    writeTextToTextbox(first_slide_shapes_dict["year"].text_frame, year)
    writeTextToTextbox(first_slide_shapes_dict["month1"].text_frame, month1)
    writeTextToTextbox(first_slide_shapes_dict["month2"].text_frame, month2)
    writeTextToTextbox(first_slide_shapes_dict["monthes"].text_frame, month1 + "-" + month2)
    writeTextToTextbox(first_slide_shapes_dict["contact"].text_frame, contact)
    writeTextToTextbox(first_slide_shapes_dict["phone"].text_frame, contact.split(' ')[1])
    


def create_date_string(month_and_day_string, year_string):
    arr = month_and_day_string.split('.')
    day = arr[0]
    month = arr[1]
    month = month.replace(" ", "")
    date_string = day + "." + month + "." + year_string
    return date_string


def processSecondSlide(slide, isFirstPresentation):
    single_event_shapes, double_event_shapes, header_shape = get_slide_shapes(slide, isFirstSlide = False)

    calendar.setfirstweekday(6)

    month = None 
    if isFirstPresentation:
        month = MetaData.firstMonthInteger
    else:
        month = MetaData.secondMonthInteger

    increaseYear = MetaData.increaseYear if not isFirstPresentation else False

    calendar.setfirstweekday(6)

    createCalendarDates(slide, single_event_shapes, double_event_shapes, month, increaseYear)

    writeTextToTextboxes(slide, single_event_shapes, double_event_shapes, header_shape, month, increaseYear, isFirstPresentation)



def clearTextboxText(text_frame):
    for i in range(len(text_frame.paragraphs)):
        para = text_frame.paragraphs[i]
        for j in range (len(para.runs)):
            para.runs[j].text=""

def get_number_of_shape(year, month, day):
    numOfDaysInMonth = calendar.monthrange(year, month)[1]
    x = np.array(calendar.monthcalendar(year, month))
    week_of_month = np.where(x==day)[0][0] # 0 is first week
    day_of_week = np.where(x == day)[1][0]+1 # 1 is Sunday
    if day_of_week > 6: 
        return -1
    first_day_of_month = np.where(x == 1)[1][0] + 1 # for removing first week that starts on Friday or Saturday
    if first_day_of_month <= 6: 
        return week_of_month*6 + day_of_week
    else:
        return (week_of_month - 1)*6 + day_of_week


def createCalendarDates(slide, singleEventShapes, doubleEventShapes, month, increaseYear = False):
    year = int(MetaData.year)
    if increaseYear:
        year += 1

    processFirstDayOffs(slide, singleEventShapes, doubleEventShapes, year, month)

    numOfDaysInMonth = calendar.monthrange(year, month)[1]
    

    for i in range(1, numOfDaysInMonth+1):
        num_of_shape = get_number_of_shape(year, month, i) - 1
        if num_of_shape < 0:
            continue
        else:
            single_event_shape = singleEventShapes[num_of_shape]
            double_event_shape = doubleEventShapes[num_of_shape]

            removeShapeOffElements(slide, single_event_shape)
            writeTextToTextbox(single_event_shape.countShape.text_frame, str(i), None, PP_ALIGN.CENTER)
            writeTextToTextbox(single_event_shape.dayShape.text_frame, hebrew_letter_of_day(get_week_day(i, month, year)), None, PP_ALIGN.CENTER) #TODO: deal with arabic

            removeShapeOffElements(slide, double_event_shape)
            writeTextToTextbox(double_event_shape.countShape.text_frame, str(i), None, PP_ALIGN.CENTER)
            writeTextToTextbox(double_event_shape.dayShape.text_frame, hebrew_letter_of_day(get_week_day(i, month, year)), None, PP_ALIGN.CENTER) #TODO: deal with arabic

    processLastDayOffs(slide, singleEventShapes, doubleEventShapes, year, month)


def removeShapeOffElements(slide, shape):
    shape.countOffShape._element.getparent().remove(shape.countOffShape._element)
    shape.dayOffShape._element.getparent().remove(shape.dayOffShape._element)
    shape.spineBgOffShape._element.getparent().remove(shape.spineBgOffShape._element)
    shape.bgOffShape._element.getparent().remove(shape.bgOffShape._element)


def processFirstDayOffs(slide, singleEventShapes, doubleEventShapes, year, month):
    prev_month = month - 1 if month != 1 else 12
    num_of_days_in_previous_month = calendar.monthrange(year, prev_month)[1]


    first_day_of_week1 = get_number_of_shape(year, month, 1) - 1


    i = first_day_of_week1 - 1
    j = num_of_days_in_previous_month
    while i >= 0:
        single_event_shape = singleEventShapes[i]
        double_event_shape = doubleEventShapes[i]

        writeTextToTextbox(single_event_shape.countOffShape.text_frame, str(j))
        writeTextToTextbox(single_event_shape.dayOffShape.text_frame, hebrew_letter_of_day(get_week_day(j, prev_month, year))) #TODO: deal with arabic
 
        treat_off_shape(slide, single_event_shape, double_event_shape)

        i -= 1
        j -= 1


def processLastDayOffs(slide, singleEventShapes, doubleEventShapes, year, month):

    next_month = month + 1 if month != 12 else 1

    num_of_days_in_month = calendar.monthrange(year, month)[1]
    last_shape_in_month_index = get_number_of_shape(year, month, num_of_days_in_month) - 1

    if last_shape_in_month_index < 0:
        last_shape_in_month_index = get_number_of_shape(year, month, num_of_days_in_month - 1) - 1

    i = last_shape_in_month_index
    j = 1

    year = year if next_month != 1 else year + 1

    while i + j < 36:
        single_event_shape = singleEventShapes[i + j]
        double_event_shape = doubleEventShapes[i+j]

        weekday = get_week_day(j, next_month, year)
        if weekday == 5:
            j += 1
            i -= 1
            continue


        writeTextToTextbox(single_event_shape.countOffShape.text_frame, str(j))
        writeTextToTextbox(single_event_shape.dayOffShape.text_frame, hebrew_letter_of_day(weekday)) #TODO: deal with arabic
        treat_off_shape(slide, single_event_shape, double_event_shape)

        j += 1



def treat_off_shape(slide, single_event_shape, double_event_shape):
    double_event_shape.shape._element.getparent().remove(double_event_shape.shape._element)
    for shape in single_event_shape.shape.shapes:
        if shape.name != "COUNT OFF" and shape.name != "DAY OFF" and shape.name != "SPINE BG OFF" and shape.name != "BG OFF":
            shape._element.getparent().remove(shape._element)


def get_week_day(day, month, year):
    date = datetime.datetime(year, month, day)
    return date.weekday()


def hebrew_letter_of_day(weekday_int):
    if weekday_int == 0:
        return "ב"
    elif weekday_int == 1:
        return "ג"
    elif weekday_int == 2:
        return "ד"
    elif weekday_int == 3:
        return "ה"
    elif weekday_int == 4:
        return "ו"
    elif weekday_int == 5:
        return ""
    elif weekday_int == 6:
        return "א"


def writeTextToTextboxes(slide, singleEventShapes, doubleEventShapes, headerShape, month, increaseYear, isFirstPresentation):
    year = int(MetaData.year)
    if increaseYear:
        year += 1

    numOfDaysInMonth = int(calendar.monthrange(year, month)[1])

    for i in range(1, numOfDaysInMonth + 1):
        shape_index = get_number_of_shape(year, month, i) - 1

        if shape_index >= 0:

            single_event_shape = singleEventShapes[shape_index]
            double_event_shepe = doubleEventShapes[shape_index]

            shape_day_in_month = int(single_event_shape.countShape.text_frame.paragraphs[0].runs[0].text)

            type = get_shape_type(shape_day_in_month, month)

            if type == ShapeType.SINGLE:
                treatSingleShape(slide, single_event_shape, double_event_shepe, shape_day_in_month, month)
            elif type == ShapeType.DOUBLE:
                treatDoubleShape(slide, single_event_shape, double_event_shepe, shape_day_in_month, month)
            elif type == ShapeType.PIC:
                treatPicShape(slide, single_event_shape, double_event_shepe)
    
    monthString = MetaData.firstMonthName if isFirstPresentation else MetaData.secondMonthName


    writeTextToTextbox(headerShape["zone"].text_frame, MetaData.zone)
    writeTextToTextbox(headerShape["month"].text_frame, monthString)
    writeTextToTextbox(headerShape["year"].text_frame, MetaData.year)


def get_shape_type(day_in_month, month):
    excelEventsDict = excelEventsDictionary[month]
    if day_in_month in excelEventsDict:
        excelEventArray = excelEventsDict[day_in_month]
        if len(excelEventArray) == 1:
            return ShapeType.SINGLE
        elif len(excelEventArray) == 2:
            return ShapeType.DOUBLE
    else:
        return ShapeType.PIC


def treatSingleShape(slide, single_event_shape, double_event_shape, shape_day_in_month, month):
    excelEvents = excelEventsDictionary[month]
    excelEvent = excelEvents[shape_day_in_month][0]

    double_event_shape.shape._element.getparent().remove(double_event_shape.shape._element)
    treatTags(slide, excelEvent, single_event_shape, None, None) #TODO: implement me

    if excelEvent.eventType is EventType.REGULAR:
        single_event_shape.bgPicShape._element.getparent().remove(single_event_shape.bgPicShape._element) #TODO: check how to properly remove pictures
        single_event_shape.bgHighlightShape._element.getparent().remove(single_event_shape.bgHighlightShape._element)
    elif excelEvent.eventType is EventType.WIDE:
        single_event_shape.bgPicShape._element.getparent().remove(single_event_shape.bgPicShape._element) #TODO: check how to properly remove pictures
        single_event_shape.bgShape._element.getparent().remove(single_event_shape.bgShape._element)

    writeTextToTextbox(single_event_shape.titleShape.text_frame, excelEvent.hour + " - " + excelEvent.title, link = excelEvent.link) 
    writeTextToTextbox(single_event_shape.locationShape.text_frame, excelEvent.location)
    writeTextToTextbox(single_event_shape.priceShape.text_frame, excelEvent.price)


def treatDoubleShape(slide, single_event_shape, double_event_shape, shape_day_in_month, month):
    excelEvents = excelEventsDictionary[month]
    firstExcelEvent = excelEvents[shape_day_in_month][0]
    secondExcelEvent = excelEvents[shape_day_in_month][1]

    single_event_shape.shape._element.getparent().remove(single_event_shape.shape._element) 

    treatTags(slide, firstExcelEvent, single_event_shape, secondExcelEvent, double_event_shape) #TODO: implement me

    double_event_shape.bgPicShape._element.getparent().remove(double_event_shape.bgPicShape._element)

    writeTextToTextbox(double_event_shape.titleShape1.text_frame, firstExcelEvent.hour + " - " + firstExcelEvent.title, link = firstExcelEvent.link)
    writeTextToTextbox(double_event_shape.locationShape1.text_frame, firstExcelEvent.location)
    writeTextToTextbox(double_event_shape.priceShape1.text_frame, firstExcelEvent.price)
    writeTextToTextbox(double_event_shape.titleShape2.text_frame, secondExcelEvent.hour + " - " + secondExcelEvent.title, link = secondExcelEvent.link)
    writeTextToTextbox(double_event_shape.locationShape2.text_frame, secondExcelEvent.location)
    writeTextToTextbox(double_event_shape.priceShape2.text_frame, secondExcelEvent.price)


def treatPicShape(slide, single_event_shape, double_event_shape):
    double_event_shape.shape._element.getparent().remove(double_event_shape.shape._element)
    single_event_shape.bgShape._element.getparent().remove(single_event_shape.bgShape._element)
    single_event_shape.bgHighlightShape._element.getparent().remove(single_event_shape.bgHighlightShape._element)
    single_event_shape.locationShape._element.getparent().remove(single_event_shape.locationShape._element)
    single_event_shape.priceShape._element.getparent().remove(single_event_shape.priceShape._element)
    treatTags(slide, None, single_event_shape, None, None)

    image_path = get_random_image()
    group_left = single_event_shape.shape.left
    group_top = single_event_shape.shape.top

    width = 1767747
    height = single_event_shape.bgPicShape.height

    slide.shapes.add_picture(image_path, group_left, group_top, width, height)
    single_event_shape.bgPicShape._element.getparent().remove(single_event_shape.bgPicShape._element)


def get_random_image():
    directory = MetaData.imagesDirectory
    all_images = MetaData.allImages
    used_images = MetaData.usedImages

    if len(used_images) == len(all_images):
        used_images.clear()
        MetaData.used_images.clear()

    unused_images = set(all_images) - used_images
    image = random.choice(list(unused_images))
    MetaData.usedImages.add(image)

    return os.path.join(directory, image)



def writeTextToTextbox(shape_text_frame, text, link = None, aligment = None):
    clearTextboxText(shape_text_frame)
    text_frame_paragraph = shape_text_frame.paragraphs[0]
    text_frame_paragraph.alignment = aligment if aligment != None else PP_ALIGN.RIGHT
    run = text_frame_paragraph.runs[0]
    run.text = text
    if link != None:
        run.hyperlink.address = link


def treatTags(slide, excelEvent, single_event_shape, secondExcelEvent, double_event_shape):
    event_shape = single_event_shape
    excelEventCommunity = excelEvent.community if excelEvent != None else None
    isDoubleShape = False
    secondExcelEventCommunity = None
    if secondExcelEvent != None:
        secondExcelEventCommunity = secondExcelEvent.community
        event_shape = double_event_shape
        isDoubleShape = True

    excelCommunitiesStrings = ExcelCommunitiesStrings()
    excelCommunitiesStringsArray = excelCommunitiesStrings.excelCommunitiesStringsArray
    communitiesStringsArray = excelCommunitiesStrings.communitiesStringsArray
    allCommunities = [e for e in Community]


    for i in range(len(allCommunities)):
        if allCommunities[i] is not excelEventCommunity:
            commString = communitiesStringsArray[i]
            suffix = "1" if secondExcelEvent != None else ""
            attrToRemove = getattr(event_shape, "tag" + commString + suffix)
            if attrToRemove != None:
                attrToRemove._element.getparent().remove(attrToRemove._element)

        if isDoubleShape:
            if allCommunities[i] is not secondExcelEventCommunity:
                commString = communitiesStringsArray[i]
                suffix = "2"
                attrToRemove = getattr(event_shape, "tag" + commString + suffix)
                if attrToRemove != None:
                    attrToRemove._element.getparent().remove(attrToRemove._element)


def removeAllTags(event_shape):
    excelCommunitiesStrings = ExcelCommunitiesStrings()
    allCommunities = [e for e in Community]
    communitiesStringsArray = excelCommunitiesStrings.communitiesStringsArray

    for i in range(len(allCommunities)):
        commString = communitiesStringsArray[i]
        attrToRemove = getattr(event_shape, "tag" + commString)
        attrToRemove._element.getparent().remove(attrToRemove._element)



def main():
    try:
        MetaData()
        create_program_GUI()
    except Exception as e:
        if str(e)=="קובץ ה-cache אינו תקין":
            easygui.msgbox("שגיאה: \nאאי אפשר לקרוא נתונים מוזנים אחרונים. הקובץ \n\"cache.txt\"\nאינו תקין")
            easygui.msgbox(e)

if __name__ == "__main__":
    main()